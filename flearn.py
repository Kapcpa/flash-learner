import argparse
import hashlib
import json
import os
import sys
import shutil
from pathlib import Path

from pypdf import PdfReader
import docx
from pptx import Presentation
import pytesseract
from PIL import Image

from groq import Groq
from dotenv import load_dotenv

from rich.console import Console
from rich.panel import Panel
from rich.prompt import Confirm, Prompt

load_dotenv()

console = Console()

FLEARN_API_KEY = None
FLEARN_FOLDER = None
FLEARN_DEBUG = False

CONFIG_DIR = Path.home() / ".flearn"
CONFIG_FILE = CONFIG_DIR / ".env"

FLEARN_FLASHCARD_GEN_PROMPT = """
You are a precise flashcard generator. Read the provided study materials and extract the most important concepts.

Rules:
1. Ignore all introductory filler, table of contents, opinions, conversational text or trivia.
2. Focus exclusively on "What is [Concept]?", "Why is [Concept] used?", and "How does [Mechanism] work?".
3. If a page or slide contains no concrete definitions, skip it entirely.
4. Keep the "back" of the flashcard concise and factu

Output ONLY valid JSON in the following exact format:
{
    "cards": [
		{"front": "What is the definition of [Concept]?", "back": "[Strict factual definition]"}  
	]
}
"""

client: Groq | None = None


def dprint(msg: str):
	if FLEARN_DEBUG:
		console.print(f"[dim][DEBUG] {msg}[/dim]")


def print_error(msg: str):
	console.print(f"[bold red][ ERROR ][/bold red] {msg}")


def print_success(msg: str):
	console.print(f"[bold green][ SUCCESS ][/bold green] {msg}")


def print_info(msg: str):
	console.print(f"[bold yellow][ INFO ][/bold yellow] {msg}")


def setup_wizard():
	CONFIG_DIR.mkdir(parents=True, exist_ok=True)

	if not CONFIG_FILE.exists():
		console.print(Panel("Welcome to [bold cyan]flearn[/bold cyan] setup", style="blue"))

		api_key = Prompt.ask("Enter your [bold]Groq API Key[/bold]", password=True).strip()
		db_folder = Prompt.ask("Where should flashcard data be saved?", default="~/.flearn").strip()

		env_content = f"FLEARN_API_KEY={api_key}\nFLEARN_FOLDER={db_folder}\n"
		CONFIG_FILE.write_text(env_content)
		print_success(f"Configuration saved to {CONFIG_FILE}\n")

	load_dotenv(dotenv_path=CONFIG_FILE)

	global FLEARN_API_KEY, FLEARN_FOLDER
	FLEARN_API_KEY = os.environ.get("FLEARN_API_KEY")
	folder_str = os.environ.get('FLEARN_FOLDER', '~/.flearn')
	FLEARN_FOLDER = Path(folder_str).expanduser()
	FLEARN_FOLDER.mkdir(parents=True, exist_ok=True)

	if getattr(sys, 'frozen', False):
		current_path = Path(sys.executable).resolve()
		target_dir = Path.home() / ".local" / "bin"
		target_path = target_dir / "flearn"

		if current_path != target_path:
			install = Confirm.ask(f"Do you want to install flearn globally to [bold]{target_path}[/bold]?", default=True)
			if install:
				try:
					target_dir.mkdir(parents=True, exist_ok=True)
					shutil.copy2(current_path, target_path)

					target_path.chmod(0o755)
					print_success("Installation complete!")
					console.print("[dim]You can now run 'flearn' from anywhere. (You may need to restart your terminal).[/dim]")

					sys.exit(0)
				except Exception as e:
					print_error(f"Failed to install globally: {e}")


def get_client() -> Groq:
	global client

	if client is None:
		if not FLEARN_API_KEY:
			print_error("FLEARN_API_KEY is not set. Please check your .env file.")
			sys.exit(1)
		client = Groq(api_key=FLEARN_API_KEY)

	return client


def get_file_hash(filepath: Path) -> str:
	hasher = hashlib.md5()
	hasher.update(filepath.read_bytes())
	return hasher.hexdigest()


def load_database(database_path: Path) -> dict:
	if not database_path.exists():
		return {}

	try:
		return json.loads(database_path.read_text())
	except json.JSONDecodeError:
		console.print(
			f"[bold yellow][ WARNING ][/bold yellow] Database file '{database_path.name}' is corrupted. Starting fresh.")
		return {}


def get_all_cards(data: dict) -> list[dict]:
	cards_by_file = data.get("cards_by_file", {})

	if cards_by_file:
		cards = []

		for file_cards in cards_by_file.values():
			cards.extend(file_cards)

		return cards

	return data.get("cards", [])


def get_llm_flashcards(content: str) -> list[dict]:
	try:
		response = get_client().chat.completions.create(
			model="llama-3.3-70b-versatile",
			messages=[
				{"role": "system", "content": FLEARN_FLASHCARD_GEN_PROMPT},
				{"role": "user", "content": content}
			],
			response_format={"type": "json_object"},
			temperature=0.3
		)

		result = json.loads(response.choices[0].message.content)
		return result.get("cards", [])
	except Exception as e:
		print_error(f"Communicating with Groq: {e}")
		return []


def text_from_file(filepath: Path) -> str:
	ext = filepath.suffix.lower()
	text = ""

	try:
		if ext in ['.txt', '.md']:
			text = filepath.read_text(encoding='utf-8')
		elif ext == '.pdf':
			reader = PdfReader(filepath)
			for page in reader.pages:
				extracted = page.extract_text()
				if extracted:
					text += extracted + "\n"
		elif ext == '.docx':
			doc = docx.Document(filepath)
			text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
		elif ext == '.pptx':
			prs = Presentation(filepath)
			for slide in prs.slides:
				for shape in slide.shapes:
					if hasattr(shape, "text"):
						text += shape.text + "\n"
		elif ext in ['.png', '.jpg', '.jpeg']:
			img = Image.open(filepath)
			text = pytesseract.image_to_string(img)
	except Exception as e:
		print_error(f"Extracting text from {filepath.name}: {e}")

	return text


def review_cards(cards: list[dict]) -> list[dict]:
	approved_cards = []

	if not cards:
		return approved_cards

	console.print(f"\n[bold magenta]--- Reviewing {len(cards)} generated cards ---[/bold magenta]")

	for i, card in enumerate(cards, 1):
		content = f"[bold cyan]Q:[/bold cyan] {card.get('front')}\n[bold yellow]A:[/bold yellow] {card.get('back')}"
		console.print(Panel(content, title=f"Card {i}/{len(cards)}", border_style="blue", padding=(1, 2)))

		keep = Confirm.ask("Keep this card?", default=True)
		if keep:
			approved_cards.append(card)
		else:
			console.print("[dim]-> Card discarded.[/dim]\n")

	return approved_cards


def sync_files(target_dir: Path, file_states: dict, cards_by_file: dict, force_regen: bool) -> bool:
	has_changes = False
	current_files = set()

	supported_extensions = ['.txt', '.md', '.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg']

	for filepath in target_dir.glob("*"):
		filename = filepath.name

		if filepath.suffix.lower() not in supported_extensions:
			console.print(f"[bold yellow][ WARNING ][/bold yellow] Filetype of {filename} not supported. Skipping.")
			continue

		current_files.add(filename)

		try:
			current_hash = get_file_hash(filepath)
			is_modified = file_states.get(filename) != current_hash

			if force_regen or is_modified:
				if FLEARN_DEBUG:
					status = "Regenerating" if force_regen else ("New" if filename not in file_states else "Modified")
					dprint(f"{status} file detected: {filename}")

				with console.status(f"[bold cyan]Extracting & Generating cards for {filename}...[/bold cyan]", spinner="dots"):
					text = text_from_file(filepath)

					if not text:
						console.print(
							f"[bold yellow][ WARNING ][/bold yellow] No readable text found in {filename}. Skipping.")
						continue

					generated_cards = get_llm_flashcards(text)

				if generated_cards:
					cards_by_file[filename] = review_cards(generated_cards)
					file_states[filename] = current_hash
					has_changes = True
				else:
					console.print(f"[bold yellow][ WARNING ][/bold yellow] Failed to generate cards for {filename}")
			else:
				dprint(f"Unchanged, skipping: {filename}")
		except Exception as e:
			console.print(f"[bold yellow][ WARNING ][/bold yellow] Could not process {filename}: {e}")

	deleted_files = [file for file in file_states.keys() if file not in current_files]

	for deleted_file in deleted_files:
		dprint(f"Removing deleted file from state: {deleted_file}")

		file_states.pop(deleted_file, None)
		cards_by_file.pop(deleted_file, None)
		has_changes = True

	return has_changes


def gen(args, force_regen: bool = False):
	target_dir = Path(args.directory).resolve()
	group_name = target_dir.name
	database_path = FLEARN_FOLDER / f"{group_name}.json"

	dprint(f"Target Directory: {target_dir}")
	dprint(f"Group Name: {group_name}")
	dprint(f"Saving to Database Path: {database_path}")

	if not target_dir.exists() or not target_dir.is_dir():
		print_error(f"Directory '{args.directory}' does not exist.")
		return

	database_data = load_database(database_path)
	file_states = {} if force_regen else database_data.get("file_states", {})
	cards_by_file = {} if force_regen else database_data.get("cards_by_file", {})

	action_str = "Regenerating all" if force_regen else "Scanning for new or modified"
	console.print(Panel(f"{action_str} files in [bold cyan]{group_name}[/bold cyan]...", style="blue"))

	has_changes = sync_files(target_dir, file_states, cards_by_file, force_regen)

	if not has_changes:
		print_success("No new or modified files found. Everything is up to date!")
		return

	save_data = {
		"group": group_name,
		"source_dir": str(target_dir),
		"file_states": file_states,
		"cards_by_file": cards_by_file
	}

	database_path.write_text(json.dumps(save_data, indent=4))
	print_success(f"New flashcards saved to group '{group_name}'.")


def rm(args):
	group_name = args.group
	target_card_numbers = sorted(args.card_numbers, reverse=True)
	database_path = FLEARN_FOLDER / f"{group_name}.json"

	dprint(f"Removing cards {target_card_numbers} from group {group_name}")

	if not database_path.exists():
		print_error(f"No flashcards found for group '{group_name}'.")
		return

	database_data = load_database(database_path)
	cards_by_file = database_data.get("cards_by_file", {})

	total_cards_deleted = 0

	for target_card_num in target_card_numbers:
		current_index = 0
		deleted = False
		target_index_0_based = target_card_num - 1

		for filename, file_cards in cards_by_file.items():
			if current_index <= target_index_0_based < current_index + len(file_cards):
				local_index = target_index_0_based - current_index
				removed_card = file_cards.pop(local_index)

				print_success(f"Deleted Card {target_card_num}:")
				console.print(f"  [dim]Q: {removed_card.get('front')}[/dim]")
				deleted = True
				total_cards_deleted += 1
				break

			current_index += len(file_cards)

		if not deleted:
			console.print(f"[bold yellow][ WARNING ][/bold yellow] Card number {target_card_num} not found. Skipping.")

	if total_cards_deleted > 0:
		save_data = {
			"group": group_name,
			"source_dir": database_data.get("source_dir", ""),
			"file_states": database_data.get("file_states", {}),
			"cards_by_file": cards_by_file
		}
		database_path.write_text(json.dumps(save_data, indent=4))
		print_success(f"Total cards removed: {total_cards_deleted}")
	else:
		print_info("No cards were removed.")


def view(args):
	group_name = args.group
	database_path = FLEARN_FOLDER / f"{group_name}.json"

	dprint(f"Looking up Group Name: {group_name}")
	dprint(f"Reading from Database Path: {database_path}")

	if not database_path.exists():
		print_error(f"No flashcards found for group '{group_name}'.")
		print_info("Run 'flearn ls' to see avaliable groups.")
		return

	database_data = load_database(database_path)
	cards = get_all_cards(database_data)

	if not cards:
		print_info(f"No cards in group '{group_name}'")
		return

	console.print(Panel(f"Deck: [bold cyan]{group_name}[/bold cyan] ({len(cards)} cards)", style="blue"))

	for i, card in enumerate(cards, 1):
		content = f"[bold cyan]Q:[/bold cyan] {card.get('front')}\n[bold yellow]A:[/bold yellow] {card.get('back')}"
		console.print(Panel(content, title=f"Card {i}", border_style="magenta", padding=(1, 2)))


def ls():
	dprint(f"Scanning Database Folder: {FLEARN_FOLDER}")

	if not FLEARN_FOLDER.exists():
		print_info("No groups found (database folder is empty).")
		return

	groups = list(FLEARN_FOLDER.glob("*.json"))
	if not groups:
		print_info("No groups found.")
		return

	console.print("\n[bold]Available flashcard groups:[/bold]")
	for file in groups:
		console.print(f"  [bold cyan]•[/bold cyan] {file.stem}")
	console.print()


def study(args):
	group_name = args.group
	database_path = FLEARN_FOLDER / f"{group_name}.json"

	dprint(f"Loading Study Session for Group: {group_name}")

	if not database_path.exists():
		print_error(f"No flashcards found for group '{group_name}'.")
		print_info("Run 'flearn ls' to see available groups.")
		return

	database_data = load_database(database_path)
	cards = get_all_cards(database_data)

	if not cards:
		print_info(f"No cards in group '{group_name}'")
		return

	for i, card in enumerate(cards, 1):
		console.clear()
		console.print(Panel(f"Study Session: [bold cyan]{group_name}[/bold cyan]", style="blue"))
		console.print("[dim]Press [Enter] to reveal answers. Type 'q' to quit.[/dim]\n")

		question_content = f"[bold cyan]Q:[/bold cyan] {card.get('front')}"
		console.print(Panel(question_content, title=f"Card {i} of {len(cards)}", border_style="blue", padding=(1, 2)))

		user_input = Prompt.ask("\n[dim]Reveal answer...[/dim]").strip().lower()
		if user_input == 'q':
			break

		console.clear()
		console.print(Panel(f"Study Session: [bold cyan]{group_name}[/bold cyan]", style="blue"))
		console.print("[dim]Press [Enter] to reveal answers. Type 'q' to quit.[/dim]\n")

		full_content = f"[bold cyan]Q:[/bold cyan] {card.get('front')}\n[bold yellow]A:[/bold yellow] {card.get('back')}"
		console.print(Panel(full_content, title=f"Card {i} of {len(cards)}", border_style="green", padding=(1, 2)))

		if i < len(cards):
			user_input = Prompt.ask("\n[dim]Next card...[/dim]").strip().lower()
			if user_input == 'q':
				break

	console.print()
	print_success("Study session complete.")


def main():
	setup_wizard()

	parser = argparse.ArgumentParser(prog="flearn", description="CLI AI-supported tool for quickly creating flashcards from materials you put in.")
	parser.add_argument("--debug", action="store_true", help="enable debug output")

	subparsers = parser.add_subparsers(dest="command", title="commands")

	# flearn gen <directory>
	parser_gen = subparsers.add_parser("gen", help="generates flashcards from new data in a given directory")
	parser_gen.add_argument("directory", type=str, help="target directory")

	# flearn regen <directory>
	parser_regen = subparsers.add_parser("regen", help="regenerates flashcards from all the data in a given directory")
	parser_regen.add_argument("directory", type=str, help="target directory")

	# flearn rm <group> <card_numbers...>
	parser_rm = subparsers.add_parser("rm", help="removes specific flashcards by their numbers")
	parser_rm.add_argument("group", type=str, help="name of the flashcard group")
	parser_rm.add_argument("card_numbers", type=int, nargs="+", help="the numbers of the cards to discard (as shown in view, separated by spaces)")

	# flearn view <group>
	parser_view = subparsers.add_parser("view", help="views flashcards from a given group")
	parser_view.add_argument("group", type=str, help="name of the flashcard group to view")

	# flearn ls
	_ = subparsers.add_parser("ls", help="lists available flashcard groups")

	# flearn study <group>
	parser_study = subparsers.add_parser("study", help="interactively study flashcards from a given group")
	parser_study.add_argument("group", type=str, help="name of the flashcard group to study")

	args = parser.parse_args()

	if args.command is None:
		parser.print_help()
		return

	global FLEARN_DEBUG
	FLEARN_DEBUG = args.debug

	try:
		if args.command == "gen":
			gen(args, False)
		elif args.command == "regen":
			gen(args, True)
		elif args.command == "rm":
			rm(args)
		elif args.command == "view":
			view(args)
		elif args.command == "ls":
			ls()
		elif args.command == "study":
			study(args)
	except KeyboardInterrupt:
		print_info("Process interrupted by user. Exiting...")


if __name__ == "__main__":
	main()